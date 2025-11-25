import os
from pathlib import Path
from io import BytesIO

# Office 处理库
try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
from .base import BaseDetector  # 继承自你的 BaseDetector

def generate_image_dir(filepath, base_dir="/tmp/service/images"):
    """
    根据源文件生成图片保存目录
    """
    file_stem = Path(filepath).stem
    save_dir = os.path.join(base_dir, file_stem)
    os.makedirs(save_dir, exist_ok=True)
    return save_dir



class OfficeDetector(BaseDetector):
    """增强版 Office 文件检测器"""

    def extract_text(self, filepath, ftype=None):
        """
        提取 Office 文件文本，尽量完整
        """
        #ftype = ftype or detect_office_type(filepath)
        text = ""

        if ftype in ("docx", "doc") and Document:
            doc = Document(filepath)
            # 段落
            text += "\n".join([p.text for p in doc.paragraphs if p.text]) + "\n"
            # 表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = " ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"
            # 页眉页脚
            # for section in doc.sections:
            #     if section.header:
            #         text += section.header.text + "\n"
            #     if section.footer:
            #         text += section.footer.text + "\n"

        elif ftype in ("xlsx", "xls") and openpyxl:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(c) for c in row if c is not None])
                    if row_text:
                        text += row_text + "\n"

        elif ftype in ("pptx", "ppt") and Presentation:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    # 文本框
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    # 表格
                    if hasattr(shape, "table"):
                        table = shape.table
                        for r in table.rows:
                            row_text = " ".join([c.text for c in r.cells])
                            text += row_text + "\n"
                    # 占位符
                    if hasattr(shape, "placeholders"):
                        for ph in shape.placeholders:
                            if hasattr(ph, "text") and ph.text:
                                text += ph.text + "\n"
        else:
            raise ValueError(f"无法处理该文件类型或缺少依赖库: {filepath}")

        return text

    def extract_images(self, filepath, base_dir="/tmp/service/images", ftype=None):
        """
        保存 Office 内嵌图片到 base_dir/源文件名/ 
        返回图片路径列表
        """
        #ftype = ftype or detect_office_type(filepath)
        save_dir = generate_image_dir(filepath, base_dir)
        imgs = []

        if ftype in ("docx", "doc") and Document:
            doc = Document(filepath)
            for idx, rel in enumerate(doc.part.rels.values()):
                if "image" in rel.reltype:
                    blob = rel.target_part.blob
                    ext = rel.target_part.content_type.split("/")[-1]
                    img_path = os.path.join(save_dir, f"{Path(filepath).stem}_img_{idx}.{ext}")
                    with open(img_path, "wb") as f:
                        f.write(blob)
                    imgs.append(img_path)

        elif ftype in ("pptx", "ppt") and Presentation:
            prs = Presentation(filepath)
            idx = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        blob = shape.image.blob
                        ext = shape.image.content_type.split("/")[-1]
                        img_path = os.path.join(save_dir, f"{Path(filepath).stem}_img_{idx}.{ext}")
                        with open(img_path, "wb") as f:
                            f.write(blob)
                        imgs.append(img_path)
                        idx += 1

        elif ftype in ("xlsx", "xls") and openpyxl:
            wb = openpyxl.load_workbook(filepath, read_only=True)
            idx = 0
            for sheet in wb.worksheets:
                for img in getattr(sheet, "_images", []):
                    blob = getattr(img, "ref", None)
                    if blob:
                        img_path = os.path.join(save_dir, f"{Path(filepath).stem}_img_{idx}.png")
                        with open(img_path, "wb") as f:
                            f.write(blob)
                        imgs.append(img_path)
                        idx += 1

        else:
            raise ValueError(f"无法提取图片: {filepath}")

        return imgs